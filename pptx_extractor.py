# pptx_extractor.py
# --------------------------------------------------------------------------
# This file reads text out of a PowerPoint (.pptx) file, one slide at a time.
#
# It collects text from:
#   - text boxes and titles
#   - tables
#   - grouped shapes / SmartArt (by looking inside groups)
#   - PICTURES on the slide (by running OCR on them)
#
# IMPORTANT: python-pptx can only read the NEWER .pptx format, not old .ppt.
# --------------------------------------------------------------------------

import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from ocr import ocr_image_bytes


def _collect_text_from_shapes(shapes, lines: list) -> None:
    """
    Go through a group of shapes and collect their text into `lines`.
    Handles groups (by going inside them) and pictures (by OCR).
    """
    for shape in shapes:

        # If this shape is a group, dig into the shapes inside it.
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _collect_text_from_shapes(shape.shapes, lines)
            continue

        # Normal text boxes / titles.
        if shape.has_text_frame:
            text = shape.text_frame.text
            if text and text.strip():
                lines.append(text.strip())

        # Tables.
        if shape.has_table:
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                row_text = " | ".join(c for c in cells if c)
                if row_text:
                    lines.append(row_text)

        # Pictures: run OCR so we can read text that lives inside images.
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image_bytes = shape.image.blob
                ocr_text = ocr_image_bytes(image_bytes)
                if ocr_text:
                    lines.append(ocr_text)
            except Exception:
                # If a picture can't be read, just skip it quietly.
                pass


def extract_text_from_pptx(file_bytes: bytes) -> list:
    """
    Read every slide in a .pptx file and pull out its text.

    Returns:
        A list of dictionaries, one per slide, like:
        [{"number": 1, "text": "..."}, {"number": 2, "text": "..."}]
    """
    presentation = Presentation(io.BytesIO(file_bytes))

    slides_data = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        lines = []
        _collect_text_from_shapes(slide.shapes, lines)
        slides_data.append({
            "number": slide_number,
            "text": "\n".join(lines),
        })

    return slides_data