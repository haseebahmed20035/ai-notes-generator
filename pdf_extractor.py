import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ocr import ocr_image_bytes


MIN_IMAGE_WIDTH = 200
MIN_IMAGE_HEIGHT = 200


def _collect_text_from_shapes(shapes, lines: list, enable_ocr: bool = True) -> None:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _collect_text_from_shapes(shape.shapes, lines, enable_ocr)
            continue

        if shape.has_text_frame:
            text = shape.text_frame.text

            if text and text.strip():
                lines.append(text.strip())

        if shape.has_table:
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                row_text = " | ".join(cell for cell in cells if cell)

                if row_text:
                    lines.append(row_text)

        if enable_ocr and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image

                if image.size[0] < MIN_IMAGE_WIDTH or image.size[1] < MIN_IMAGE_HEIGHT:
                    continue

                ocr_text = ocr_image_bytes(image.blob)

                if ocr_text:
                    lines.append("[Image OCR]\n" + ocr_text)

            except Exception:
                pass


def extract_text_from_pptx(file_bytes: bytes, enable_ocr: bool = True) -> list:
    presentation = Presentation(io.BytesIO(file_bytes))
    slides_data = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        lines = []

        _collect_text_from_shapes(
            slide.shapes,
            lines,
            enable_ocr=enable_ocr,
        )

        slides_data.append(
            {
                "number": slide_number,
                "text": "\n".join(lines).strip(),
            }
        )

    return slides_data